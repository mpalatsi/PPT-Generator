from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import logging

logger = logging.getLogger(__name__)

class PowerPointService:
    """Service for creating PowerPoint presentations."""
    
    def __init__(self):
        """Initialize the PowerPoint service."""
        self.title_font_size = Pt(44)
        self.body_font_size = Pt(32)
    
    def create_presentation(self, screenshot_paths: list[str], output_path: str) -> None:
        """
        Create a PowerPoint presentation from a list of screenshot paths.
        
        Args:
            screenshot_paths: List of paths to screenshot images
            output_path: Path where the presentation should be saved
        """
        try:
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = "Web Screenshots Presentation"
            subtitle.text = f"Generated with {len(screenshot_paths)} slides"
            
            # Add content slides
            for i, screenshot_path in enumerate(screenshot_paths, 1):
                # Add new slide with title and content layout
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                # Add title
                title = slide.shapes.add_textbox(
                    Inches(0.5),
                    Inches(0.5),
                    Inches(9),
                    Inches(0.5)
                )
                title.text = f"Slide {i}"
                
                # Add screenshot
                if os.path.exists(screenshot_path):
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(5)
                    
                    slide.shapes.add_picture(
                        screenshot_path,
                        left,
                        top,
                        width=width,
                        height=height
                    )
                else:
                    logger.error(f"Screenshot not found: {screenshot_path}")
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"Presentation saved to {output_path}")
            
        except Exception as e:
            logger.error(f"Error creating presentation: {str(e)}")
            raise 